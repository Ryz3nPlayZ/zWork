"""Tool execution for the chat agent.

Each tool is:
  - declared in TOOL_SCHEMAS for the LLM
  - implemented here with a clear contract:
      yields `activity` events (for the UI)
      yields exactly one `tool_result` event at the end
"""

from __future__ import annotations

import asyncio
import contextlib
import json
import os
import re
import signal
import subprocess
import urllib.parse
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, AsyncIterator

import httpx

# ---------------- Tool schemas (provider-neutral) ----------------
from . import academic as academic_mod
from . import skills as skills_mod
from .home import memory_path
from .runtime import current_run

READ_ONLY_TOOLS = frozenset(
    {
        "read_file",
        "list_dir",
        "read_skill",
        "extract_document",
        "web_search",
        "search_papers",
        "format_citation",
    }
)

READ_ONLY_DCTL_SUBCOMMANDS = frozenset(
    {
        "snapshot",
        "tree",
        "screenshot",
        "windows",
        "apps",
        "describe",
        "list",
    }
)

DESTRUCTIVE_COMMAND_PATTERNS = (
    (re.compile(r"\brm\s+-[a-z]*r[a-z]*\b", re.I), "recursive delete"),
    (re.compile(r"\bmkfs(?:\.\w+)?\b", re.I), "format disk/filesystem"),
    (re.compile(r"\bdd\s+[^;&|]*\bof=/dev/", re.I), "raw disk write"),
    (re.compile(r"\bgit\s+push\b[^;&|]*(?:--force|-f)\b", re.I), "force push"),
    (re.compile(r"\bgit\s+reset\b[^;&|]*--hard\b", re.I), "hard reset"),
    (
        re.compile(r"\bgit\s+clean\b[^;&|]*-[a-z]*[fd][a-z]*", re.I),
        "destructive git clean",
    ),
    (re.compile(r"\bgit\s+branch\b[^;&|]*\s-D\b", re.I), "delete branch"),
    (re.compile(r"\bdrop\s+(?:table|database)\b", re.I), "database drop"),
    (re.compile(r"\b(?:shutdown|reboot|halt|poweroff)\b", re.I), "system shutdown"),
)


def _normalized_command(command: str) -> str:
    return " ".join(str(command or "").strip().split())


def _targets_zwork_backend(command: str) -> bool:
    c = _normalized_command(command).lower()
    if "8787" not in c:
        return False
    return bool(
        re.search(r"\blsof\b[^;&|]*(?::8787|-i\s*:8787)", c)
        and re.search(r"\b(?:xargs\s+)?kill(?:all)?\b", c)
    ) or bool(re.search(r"\b(?:kill|pkill|killall)\b[^;&|]*8787", c))


def _matches_destructive_command(command: str) -> tuple[bool, str]:
    if _targets_zwork_backend(command):
        return True, "kills the zWork local backend on port 8787"
    for pattern, reason in DESTRUCTIVE_COMMAND_PATTERNS:
        if pattern.search(command or ""):
            return True, reason
    return False, ""


def tool_risk(tool_name: str, params: dict[str, Any]) -> tuple[str, str]:
    if tool_name in READ_ONLY_TOOLS:
        return "safe", "read-only tool"
    if tool_name == "search_papers":
        return "safe", "searches academic databases"
    if tool_name == "spawn_agent":
        return "safe", "spawns a subagent which follows its own permission checks"
    if tool_name == "run_command":
        command = str(params.get("command") or "")
        destructive, reason = _matches_destructive_command(command)
        if destructive:
            return "destructive", reason
        return "sensitive", "runs a shell command"
    if tool_name == "write_file":
        return "sensitive", "writes or overwrites a file"
    if tool_name == "deploy_web_app":
        return "sensitive", "starts a local server"
    if tool_name == "save_memory":
        return "sensitive", "persists memory"
    if tool_name == "dctl_system":
        action = str(params.get("action") or "").strip().lower()
        if action in ("capabilities", "doctor", "list-apps", "list-windows", "list-launchable"):
            return "safe", "read-only system discovery"
        return "sensitive", "starts apps or opens resources"
    if tool_name == "dctl_ui":
        action = str(params.get("action") or "").strip().lower()
        if action in ("tree", "element", "read", "describe", "screenshot"):
            return "safe", "read-only desktop UI inspection"
        return "sensitive", "controls the desktop UI"
    if tool_name == "dctl_browser":
        action = str(params.get("action") or "").strip().lower()
        if action in ("tabs", "targets", "active-tab", "dom", "ax", "text", "selector", "actions", "selection", "caret", "snapshot"):
            return "safe", "read-only browser inspection"
        return "sensitive", "controls the browser state"
    if tool_name == "dctl_office":
        action = str(params.get("action") or "").strip().lower()
        if action in ("inspect", "read", "paragraphs", "sheets"):
            return "safe", "read-only document inspection"
        return "sensitive", "modifies documents or spreadsheets"
    if tool_name.startswith("mcp__"):
        return "sensitive", "calls an external MCP tool"
    if tool_name.startswith("composio__"):
        return "sensitive", "calls an external app action via Composio"
    return "sensitive", "tool changes or accesses external state"


def filter_tools_for_plan_mode(schemas: list[dict]) -> list[dict]:
    return [t for t in schemas if t.get("name") in READ_ONLY_TOOLS]


TOOL_SCHEMAS: list[dict] = [
    {
        "name": "write_file",
        "description": (
            "Write content to a file at the given path. Creates parent directories if needed. "
            "Overwrites existing files. Use this for creating app files, code, docs, configs."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative or absolute file path",
                },
                "content": {
                    "type": "string",
                    "description": "Full UTF-8 content of the file",
                },
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "read_file",
        "description": (
            "Read and return the UTF-8 contents of a file. Use this to inspect files you or "
            "the user just wrote/changed, or to see existing project files before editing."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative or absolute file path",
                },
            },
            "required": ["path"],
        },
    },
    {
        "name": "list_dir",
        "description": "List the immediate children of a directory. Use this to orient yourself.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Directory path (default: '.')",
                },
            },
            "required": [],
        },
    },
    {
        "name": "run_command",
        "description": (
            "Run a shell command. Set background=true for long-running servers (e.g. dev servers); "
            "the command will detach and return immediately. For foreground commands, the combined "
            "stdout+stderr is returned (120s timeout). Do NOT use this for interacting with "
            "connected apps (email, calendar, Slack, etc.) — use composio__ tools instead."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "command": {
                    "type": "string",
                    "description": "Shell command to execute",
                },
                "cwd": {
                    "type": "string",
                    "description": "Working directory (default: '.')",
                },
                "background": {
                    "type": "boolean",
                    "description": "Run detached (for servers)",
                },
            },
            "required": ["command"],
        },
    },
    {
        "name": "read_skill",
        "description": (
            "Load the full SKILL.md for an installed skill so you can follow its playbook. "
            "Pass the skill slug (e.g. 'anthropic-skills/pdf' or 'uiux-pro-max'). "
            "The system prompt lists available skills. Call this when a user task matches a skill's domain."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "slug": {
                    "type": "string",
                    "description": "Skill slug or leaf folder name",
                },
            },
            "required": ["slug"],
        },
    },
    {
        "name": "deploy_web_app",
        "description": (
            "Serve a local web app directory on http://localhost:<port>. Picks a free port "
            "(prefers 5173, 8000, 3000). Uses `python3 -m http.server` for static sites, or "
            "`npm run dev` when package.json has a dev script. Returns the URL."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "project_path": {
                    "type": "string",
                    "description": "Path to the project root",
                },
                "framework": {
                    "type": "string",
                    "description": "Framework hint (optional)",
                },
            },
            "required": ["project_path"],
        },
    },
    {
        "name": "save_memory",
        "description": (
            "Append a fact or note to persistent memory. "
            "MUST be called whenever the user asks to remember, note down, or save something. "
            "Do NOT just acknowledge the request — actually invoke this tool."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content": {
                    "type": "string",
                    "description": "The fact or note to remember",
                },
            },
            "required": ["content"],
        },
    },
    {
        "name": "extract_document",
        "description": (
            "Extract text (and where applicable, tables and metadata) from a document on disk. "
            "Supports PDF, DOCX, XLSX, PPTX, TXT and Markdown — auto-detected from the file "
            "extension. Use this before answering questions about the contents of a file the "
            "user has pointed at; do not try to read these formats with read_file."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the document"},
                "format": {
                    "type": "string",
                    "description": "Output style hint ('markdown' or 'text'); default 'markdown'",
                },
                "pages": {
                    "type": "string",
                    "description": "Optional 1-based page range for PDFs, e.g. '1-5' or '3'",
                },
            },
            "required": ["path"],
        },
    },
    {
        "name": "web_search",
        "description": (
            "Search the web/news without opening a browser. Use this for current events, "
            "recent news, or factual web lookup requests when the user wants the answer in chat. "
            "Do NOT use for academic papers (use search_papers) or for connected app actions "
            "like email or calendar (use composio__ tools)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Search query. Leave empty for top current headlines.",
                },
                "max_results": {
                    "type": "integer",
                    "description": "Maximum results to return (default 6, max 10)",
                },
            },
            "required": [],
        },
    },
    {
      "name": "dctl_system",
      "description": "General system control and discovery. Use this to list open windows, find launchable apps, or start new processes. Returns JSON describing windows, apps, or execution status.",
      "parameters": {
        "type": "object",
        "properties": {
          "action": {
            "type": "string",
            "enum": ["capabilities", "doctor", "list-apps", "list-windows", "list-launchable", "launch", "open"],
            "description": "The system action to perform. 'launch' starts an app by name; 'open' opens a URL or file path; 'list-windows' shows all visible window titles and IDs."
          },
          "target": {
            "type": "string",
            "description": "The target app name (for launch) or URL/file path (for open). Optional for other actions."
          },
          "cwd": {
            "type": "string",
            "description": "Working directory for the command (default: '.')"
          }
        },
        "required": ["action"]
      }
    },
    {
      "name": "dctl_ui",
      "description": "Native UI automation via the OS Accessibility Tree. Use this to inspect UI elements, click buttons, type into fields, or read text on the desktop. Locates elements using a boolean selector string (e.g. 'app:\"Code\" AND role:button').",
      "parameters": {
        "type": "object",
        "properties": {
          "action": {
            "type": "string",
            "enum": ["tree", "element", "read", "click", "type", "key", "focus", "scroll", "wait", "describe", "screenshot", "clipboard"],
            "description": "The UI interaction to perform. 'tree' dumps the accessibility hierarchy; 'click' triggers a button; 'type' inserts text; 'describe' identifies what is at coordinates; 'clipboard' reads or writes the system clipboard."
          },
          "selector": {
            "type": "string",
            "description": "The dctl selector query to find the element (e.g., 'app:\"Chrome\" AND name:\"Search\"'). Required for click, type, focus, read, wait, element."
          },
          "text": {
            "type": "string",
            "description": "The text to type into the element (used with action='type'), or text to write to clipboard (used with action='clipboard' and clipboard_action='write')."
          },
          "button": {
            "type": "string",
            "enum": ["left", "right", "middle"],
            "default": "left",
            "description": "Mouse button for click action. Use 'right' for context menu."
          },
          "double": {
            "type": "boolean",
            "default": False,
            "description": "Whether to double-click (used with action='click')."
          },
          "clipboard_action": {
            "type": "string",
            "enum": ["read", "write"],
            "description": "Clipboard operation. 'read' returns current clipboard text; 'write' stores the text parameter."
          },
          "combo": {
            "type": "string",
            "description": "The key combination to press (e.g., 'ctrl+c', 'win+r'). Used with action='key'."
          },
          "direction": {
            "type": "string",
            "enum": ["up", "down", "left", "right"],
            "description": "Scroll direction (used with action='scroll')."
          },
          "amount": {
            "type": "integer",
            "default": 1,
            "description": "Number of scroll increments."
          },
          "x": { "type": "integer", "description": "X coordinate for 'describe'." },
          "y": { "type": "integer", "description": "Y coordinate for 'describe'." },
          "cwd": {
            "type": "string",
            "description": "Working directory for the command (default: '.')"
          }
        },
        "required": ["action"]
      }
    },
    {
      "name": "dctl_browser",
      "description": "Deep browser automation via CDP (Chrome DevTools Protocol). Use this for complex web tasks like automating Google Docs, switching tabs, or reading the DOM. This bypasses OS accessibility and interacts directly with the browser engine.",
      "parameters": {
        "type": "object",
        "properties": {
          "action": {
            "type": "string",
            "enum": ["start", "tabs", "targets", "active-tab", "open", "activate", "close", "ax", "dom", "text", "selector", "actions", "selection", "caret", "click", "click-action", "act", "type", "press", "eval", "send", "wait-url", "wait-selector", "snapshot", "batch"],
            "description": "Browser action. Use `actions` + `click-action` for semantic clicking, `selector` for deterministic CSS diagnostics, and `snapshot` for structured page extraction."
          },
          "target": {
            "type": "string",
            "description": "The target tab/page ID or index. Required for most actions except 'start' and 'tabs'."
          },
          "selector": {
            "type": "string",
            "description": "CSS selector for DOM actions or AX selector."
          },
          "url": {
            "type": "string",
            "description": "URL to navigate to."
          },
          "text": {
            "type": "string",
            "description": "Text or key combo or json batch array or expression to pass to the action."
          },
          "expression": {
            "type": "string",
            "description": "JavaScript expression to evaluate in the page context."
          },
          "session": {
            "type": "string",
            "description": "Optional session name for persistent browser state."
          },
          "cwd": {
            "type": "string",
            "description": "Working directory for the command (default: '.')"
          }
        },
        "required": ["action"]
      }
    },
    {
      "name": "dctl_office",
      "description": "Semantic document and spreadsheet editing (Word/Excel/LibreOffice). Use this to read paragraphs, edit cells, append text, or replace content in documents without a GUI. Supports .docx and .xlsx files.",
      "parameters": {
        "type": "object",
        "properties": {
          "type": {
            "type": "string",
            "enum": ["word", "excel", "libreoffice"],
            "description": "The type of document or editor backend."
          },
          "action": {
            "type": "string",
            "enum": ["inspect", "read", "paragraphs", "append", "set-paragraph", "replace", "sheets", "write-cell", "write-range", "fill-table", "locate-cell", "fill-cell"],
            "description": "The editing action."
          },
          "path": {
            "type": "string",
            "description": "Path to the .docx or .xlsx file."
          },
          "text": {
            "type": "string",
            "description": "Text content to insert or append."
          },
          "index": {
            "type": "integer",
            "description": "Paragraph or element index."
          },
          "sheet": {
            "type": "string",
            "description": "Sheet name for Excel/Calc."
          },
          "cell": {
            "type": "string",
            "description": "Cell reference (e.g., 'A1')."
          },
          "value": {
            "type": "string",
            "description": "Value to write to a cell."
          },
          "find": { "type": "string", "description": "Search text or row label" },
          "replace": { "type": "string", "description": "Replacement text or column label" },
          "cwd": {
            "type": "string",
            "description": "Working directory for the command (default: '.')"
          }
        },
        "required": ["type", "action", "path"]
      }
    },
    {
        "name": "spawn_agent",
        "description": (
            "Spawn a subagent to work on a separate task in parallel. "
            "Use this when you have independent tasks that can be worked on concurrently, "
            "such as: reading multiple unrelated files, searching different directories, "
            "or running non-sequential commands. The subagent will run with the same model "
            "and tools available to you. Provide a clear task description and any necessary "
            "context as the initial message."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "task": {
                    "type": "string",
                    "description": "Clear description of what the subagent should do",
                },
                "context": {
                    "type": "string",
                    "description": "Additional context or instructions for the subagent (optional)",
                },
            },
            "required": ["task"],
        },
    },
    {
        "name": "search_papers",
        "description": (
            "Search academic literature across Semantic Scholar, arXiv, OpenAlex, and CrossRef. "
            "Returns ranked, de-duplicated papers with titles, authors, abstracts, years, "
            "citation counts, DOIs, journal names, and open-access PDF links. "
            "Use this when the user asks about academic research, scientific papers, "
            "literature reviews, or wants to find scholarly sources on a topic."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Research query (e.g. 'transformer attention mechanisms')",
                },
                "max_results": {
                    "type": "integer",
                    "description": "Maximum results (default 20, max 50)",
                },
                "year_min": {
                    "type": "integer",
                    "description": "Earliest publication year (optional, e.g. 2020)",
                },
                "year_max": {
                    "type": "integer",
                    "description": "Latest publication year (optional, e.g. 2024)",
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "format_citation",
        "description": (
            "Format a paper into a properly formatted citation string. "
            "Use this after search_papers to produce clean citations for papers the user wants to reference. "
            "Supports APA, MLA, and Chicago styles."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "paper": {
                    "type": "object",
                    "description": "A paper object from search_papers results (with title, authors, year, doi, journal)",
                },
                "style": {
                    "type": "string",
                    "description": "Citation style: apa, mla, or chicago (default: apa)",
                },
            },
            "required": ["paper"],
        },
    },
    {
        "name": "manage_tasks",
        "description": (
            "List, create, update, or delete user tasks in the cockpit's Kanban board. "
            "Use this whenever the user asks to add a task, change a task's status/column "
            "(inbox, todo, doing, done), update a task, list their tasks, or delete a task."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list", "create", "update", "delete"],
                    "description": "Action to perform on tasks",
                },
                "task_id": {
                    "type": "string",
                    "description": "ID of the task to update or delete (optional)",
                },
                "title": {
                    "type": "string",
                    "description": "Title of the task (required for create/update)",
                },
                "column": {
                    "type": "string",
                    "enum": ["inbox", "todo", "doing", "done"],
                    "description": "Column state for the task",
                },
                "due_date": {
                    "type": "string",
                    "description": "Due date in YYYY-MM-DD format (optional)",
                },
            },
            "required": ["action"],
        },
    },
    {
        "name": "manage_events",
        "description": (
            "List, create, or delete calendar events in the cockpit's Daily Agenda. "
            "Use this whenever the user asks to schedule a meeting, add an event to their calendar, "
            "list their events/schedule, or remove an event."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list", "create", "delete"],
                    "description": "Action to perform on calendar events",
                },
                "event_id": {
                    "type": "string",
                    "description": "ID of the event to delete (optional)",
                },
                "title": {
                    "type": "string",
                    "description": "Title of the event (required for create)",
                },
                "date": {
                    "type": "string",
                    "description": "Date in YYYY-MM-DD format (required for create)",
                },
                "start_time": {
                    "type": "string",
                    "description": "Start time in HH:MM format (optional, e.g. '09:00')",
                },
                "end_time": {
                    "type": "string",
                    "description": "End time in HH:MM format (optional, e.g. '10:00')",
                },
            },
            "required": ["action"],
        },
    },
]


# ---------------- Dispatcher ----------------


def _friendly_error(err: Exception, context: str = "") -> str:
    """Translate common exceptions into user-actionable messages."""
    msg = str(err)

    if isinstance(err, FileNotFoundError):
        return f"File not found: {msg}. Check the path and try again."
    if isinstance(err, NotADirectoryError):
        return f"Not a directory: {msg}. Specify a directory path instead."
    if isinstance(err, IsADirectoryError):
        return f"Is a directory, not a file: {msg}. Specify a file path instead."
    if isinstance(err, PermissionError):
        return f"Permission denied: {msg}. Check file permissions and try again."
    if isinstance(err, UnicodeDecodeError):
        return f"Cannot read as text: {msg}. The file may be binary or use an unsupported encoding."
    if isinstance(err, subprocess.TimeoutExpired):
        return f"Command timed out after {err.timeout}s. Try breaking the work into smaller steps."
    if isinstance(err, asyncio.TimeoutError):
        return "Command timed out. Try breaking the work into smaller steps."
    if isinstance(err, asyncio.CancelledError):
        return "Operation cancelled."
    if isinstance(err, OSError):
        return f"System error: {msg}. {context}Check that the path and permissions are correct."

    return f"{msg}. {context}" if context else msg


async def execute_tool(tool_name: str, params: dict[str, Any]) -> AsyncIterator[dict]:
    tool_id = f"tool_{tool_name}_{id(params)}"
    run = current_run()
    if run is not None:
        run.next_tool_call()
        run.log("tool_started", tool_name=tool_name, params=params)

    if tool_name.startswith("mcp__"):
        label = f"MCP {tool_name}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "tool",
            "done": False,
        }
        try:
            from .mcp import get_manager

            result = await get_manager().call_tool(tool_name, params)
            text = _format_mcp_result(result)
            ok = not bool(result.get("isError"))
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "tool",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=ok, output=text)
            yield {"type": "tool_result", "tool": tool_name, "ok": ok, "message": text}
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "tool",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name.startswith("composio__"):
        label = f"App: {tool_name[len('composio__') :].replace('_', ' ').title()}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "plug",
            "done": False,
        }
        try:
            from .composio import get_manager

            result = await get_manager().call_tool(tool_name, params)
            text = _format_mcp_result(result)
            ok = not bool(result.get("isError"))
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "plug",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=ok, output=text)
            yield {"type": "tool_result", "tool": tool_name, "ok": ok, "message": text}
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "plug",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "write_file":
        path = params.get("path", "")
        content = params.get("content", "")
        label = f"Write {_short_path(path)}"
        icon = _icon_for_path(path)
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": icon,
            "done": False,
        }
        try:
            await asyncio.to_thread(_write_file, path, content)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": icon,
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Wrote {len(content)} chars to {path}",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": f"Wrote {len(content)} chars to {path}",
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": icon,
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e, "Try a different path. "),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Try a different path. "),
            }
        return

    if tool_name == "read_file":
        path = params.get("path", "")
        label = f"Read {_short_path(path)}"
        icon = _icon_for_path(path)
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": icon,
            "done": False,
        }
        try:
            text = await asyncio.to_thread(_read_file, path)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": icon,
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Read {len(text)} chars from {path}",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": text,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": icon,
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e, "Try a different path. "),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Try a different path. "),
            }
        return

    if tool_name == "manage_tasks":
        action = params.get("action", "")
        task_id = params.get("task_id")
        title = params.get("title", "")
        column = params.get("column", "inbox")
        due_date = params.get("due_date")

        label = f"Tasks: {action.title()}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "list",
            "done": False,
        }
        try:
            from . import taskstore
            if action == "list":
                tasks = taskstore.get_tasks()
                res = json.dumps([taskstore.asdict(t) for t in tasks], indent=2)
            elif action == "create":
                if not title:
                    raise ValueError("Title is required to create a task")
                t = taskstore.save_task(title, column, due_date)
                res = f"Created task: {t.title} (ID: {t.id}) in column '{t.column}'"
            elif action == "update":
                if not task_id:
                    raise ValueError("Task ID is required to update a task")
                t = taskstore.save_task(title, column, due_date, task_id)
                res = f"Updated task: {t.title} (ID: {t.id}) in column '{t.column}'"
            elif action == "delete":
                if not task_id:
                    raise ValueError("Task ID is required to delete a task")
                ok = taskstore.delete_task(task_id)
                res = f"Deleted task with ID: {task_id}" if ok else f"Task with ID {task_id} not found"
            else:
                raise ValueError(f"Unknown action: {action}")

            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "list",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=True, output=res)
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": res,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "list",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=False, output=str(e))
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": f"Error managing tasks: {str(e)}",
            }
        return

    if tool_name == "manage_events":
        action = params.get("action", "")
        event_id = params.get("event_id")
        title = params.get("title", "")
        date = params.get("date", "")
        start_time = params.get("start_time")
        end_time = params.get("end_time")

        label = f"Events: {action.title()}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "calendar",
            "done": False,
        }
        try:
            from . import taskstore
            if action == "list":
                events = taskstore.get_events()
                res = json.dumps([taskstore.asdict(e) for e in events], indent=2)
            elif action == "create":
                if not title or not date:
                    raise ValueError("Title and Date are required to create an event")
                e = taskstore.save_event(title, date, start_time, end_time)
                res = f"Created event: {e.title} on {e.date} (ID: {e.id})"
            elif action == "delete":
                if not event_id:
                    raise ValueError("Event ID is required to delete an event")
                ok = taskstore.delete_event(event_id)
                res = f"Deleted event with ID: {event_id}" if ok else f"Event with ID {event_id} not found"
            else:
                raise ValueError(f"Unknown action: {action}")

            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "calendar",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=True, output=res)
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": res,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "calendar",
                "done": True,
            }
            if run is not None:
                run.log("tool_finished", tool_name=tool_name, ok=False, output=str(e))
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": f"Error managing events: {str(e)}",
            }
        return

    if tool_name == "list_dir":
        path = params.get("path", ".")
        label = f"List {_short_path(path)}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "folder",
            "done": False,
        }
        try:
            listing = await asyncio.to_thread(_list_dir, path)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "folder",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Listed {path}",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": listing,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "folder",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e, "Check that the path is a directory. "),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Check that the path is a directory. "),
            }
        return

    if tool_name == "run_command":
        command = params.get("command", "")
        cwd = params.get("cwd", ".")
        background = bool(params.get("background", False))
        short = command[:60] + ("…" if len(command) > 60 else "")
        label = f"Run: {short}" + (" (bg)" if background else "")
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "command",
            "done": False,
        }
        try:
            _ensure_command_allowed(command)
            if background:
                pid = await asyncio.to_thread(_run_background, command, cwd)
                yield {
                    "type": "activity",
                    "id": tool_id,
                    "label": label,
                    "icon": "command",
                    "done": True,
                }
                if run is not None:
                    run.log(
                        "tool_finished",
                        tool_name=tool_name,
                        ok=True,
                        output=f"Started background process (pid={pid})",
                    )
                yield {
                    "type": "tool_result",
                    "tool": tool_name,
                    "ok": True,
                    "message": f"Started background process (pid={pid})",
                }
            else:
                yield {"type": "status", "text": f"Running command: {short}"}
                result = await _run_command(command, cwd)
                yield {
                    "type": "activity",
                    "id": tool_id,
                    "label": label,
                    "icon": "command",
                    "done": True,
                }
                if run is not None:
                    run.log(
                        "tool_finished",
                        tool_name=tool_name,
                        ok=result["ok"],
                        output=result["output"] or f"exit {result['returncode']}",
                    )
                yield {
                    "type": "tool_result",
                    "tool": tool_name,
                    "ok": result["ok"],
                    "message": result["output"] or (f"exit {result['returncode']}"),
                }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "command",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "web_search":
        query = str(params.get("query") or "").strip()
        max_results = int(params.get("max_results") or 6)
        label = f"Search web: {query[:50]}" if query else "Search current headlines"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "search",
            "done": False,
        }
        try:
            text = await asyncio.to_thread(_web_search, query, max_results)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "search",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished", tool_name=tool_name, ok=True, output=text[:1000]
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": text,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "search",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "read_skill":
        slug = params.get("slug", "")
        label = f"Read skill {slug}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "file",
            "done": False,
        }
        try:
            text = await asyncio.to_thread(skills_mod.read_skill, slug)
            if text is None:
                available = ", ".join(s.slug for s in skills_mod.list_skills()[:10])
                msg = f"No skill named '{slug}'. Try one of: {available}"
                yield {
                    "type": "activity",
                    "id": tool_id,
                    "label": f"Failed: {label}",
                    "icon": "file",
                    "done": True,
                }
                if run is not None:
                    run.log("tool_finished", tool_name=tool_name, ok=False, output=msg)
                yield {
                    "type": "tool_result",
                    "tool": tool_name,
                    "ok": False,
                    "message": msg,
                }
            else:
                # Cap to keep context sane.
                if len(text) > 80_000:
                    text = text[:80_000] + "\n…[truncated]"
                yield {
                    "type": "activity",
                    "id": tool_id,
                    "label": label,
                    "icon": "file",
                    "done": True,
                }
                if run is not None:
                    run.log(
                        "tool_finished",
                        tool_name=tool_name,
                        ok=True,
                        output=f"Loaded skill {slug}",
                    )
                yield {
                    "type": "tool_result",
                    "tool": tool_name,
                    "ok": True,
                    "message": text,
                }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e, "Check the skill slug spelling. "),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Check the skill slug spelling. "),
            }
        return

    if tool_name == "deploy_web_app":
        project_path = params.get("project_path", ".")
        framework = params.get("framework", "")
        label = f"Serve {_short_path(project_path)}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "deploy",
            "done": False,
        }
        try:
            result = await asyncio.to_thread(_deploy_web_app, project_path, framework)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "deploy",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=result["ok"],
                    output=result["message"],
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": result["ok"],
                "message": result["message"],
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "deploy",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "save_memory":
        content = params.get("content", "")
        label = "Save memory"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "file",
            "done": False,
        }
        try:
            await asyncio.to_thread(_save_memory, content)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output="Saved to memory.",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": "Saved to memory.",
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "extract_document":
        path = params.get("path", "")
        fmt = params.get("format", "markdown") or "markdown"
        pages = params.get("pages")
        label = f"Extract {_short_path(path)}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "file",
            "done": False,
        }
        try:
            result = await asyncio.to_thread(_extract_document, path, fmt, pages)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Extracted document {path}",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": json.dumps(result, ensure_ascii=False),
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e, "Check the file path and extension. "),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Check the file path and extension. "),
            }
        return

    if tool_name in ("dctl_system", "dctl_ui", "dctl_browser", "dctl_office"):
        cwd = params.get("cwd", ".")
        try:
            if tool_name == "dctl_system":
                subcommand, args = _map_dctl_system(params)
            elif tool_name == "dctl_ui":
                subcommand, args = _map_dctl_ui(params)
            elif tool_name == "dctl_browser":
                subcommand, args = _map_dctl_browser(params)
            else: # dctl_office
                subcommand, args = _map_dctl_office(params)
        except Exception as e:
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": f"Parameter mapping error: {str(e)}",
            }
            return

        full = " ".join([subcommand, *args]).strip()
        label = f"dctl {full}" if full else "dctl"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "window",
            "done": False,
        }
        try:
            result = await _run_dctl(subcommand, args, cwd)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "window",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=result["ok"],
                    output=result["output"]
                    or (f"exit {result['returncode']}" if not result["ok"] else "ok"),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": result["ok"],
                "message": result["output"]
                or (f"exit {result['returncode']}" if not result["ok"] else "ok"),
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "window",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(
                        e, "Check that dctl is installed and the subcommand is valid. "
                    ),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(
                    e, "Check that dctl is installed and the subcommand is valid. "
                ),
            }
        return

    if tool_name == "spawn_agent":
        task = params.get("task", "")
        context = params.get("context", "")
        if not task:
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": "spawn_agent requires a 'task' parameter describing what to do.",
            }
            return

        label = f"Agent: {task[:40]}{'...' if len(task) > 40 else ''}"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "bot",
            "done": False,
        }

        try:
            from .subagent import spawn_agent

            # Build messages for the subagent
            sub_messages = []
            if context:
                sub_messages.append({"role": "user", "content": f"Context: {context}"})
            sub_messages.append({"role": "user", "content": task})

            # Get the current model from the run context
            model_id = run.requested_model_id if run else "claude-sonnet-4-5-20250929"

            # Stream the subagent
            full_result = []
            async for evt in spawn_agent(task, sub_messages, model_id, plan_mode=False):
                if evt.get("type") == "subagent_delta":
                    full_result.append(evt.get("text", ""))
                elif evt.get("type") == "subagent_done":
                    if evt.get("error"):
                        yield {
                            "type": "activity",
                            "id": tool_id,
                            "label": f"Failed: {label}",
                            "icon": "bot",
                            "done": True,
                        }
                        yield {
                            "type": "tool_result",
                            "tool": tool_name,
                            "ok": False,
                            "message": f"Subagent failed: {evt.get('error')}",
                        }
                        return

            result = "".join(full_result)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "bot",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Subagent completed with {len(result)} chars",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": result,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "bot",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if tool_name == "search_papers":
        query = str(params.get("query") or "").strip()
        max_results = min(int(params.get("max_results") or 20), 50)
        year_min = params.get("year_min")
        year_max = params.get("year_max")
        if year_min is not None:
            year_min = int(year_min)
        if year_max is not None:
            year_max = int(year_max)
        label = f"Search papers: {query[:50]}" if query else "Search papers"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "search",
            "done": False,
        }
        try:
            papers = await academic_mod.search_academic_literature(
                query, max_results=max_results, year_min=year_min, year_max=year_max
            )
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "search",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=True,
                    output=f"Found {len(papers)} papers for '{query}'",
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": json.dumps(papers, ensure_ascii=False),
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "search",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e, "Check the query and try again. "),
            }
        return

    if tool_name == "format_citation":
        paper_json = params.get("paper", {})
        style = str(params.get("style") or "apa").strip().lower()
        if style not in ("apa", "mla", "chicago"):
            style = "apa"
        label = f"Format citation ({style})"
        yield {
            "type": "activity",
            "id": tool_id,
            "label": label,
            "icon": "file",
            "done": False,
        }
        try:
            citation = academic_mod.format_citation(paper_json, style)
            yield {
                "type": "activity",
                "id": tool_id,
                "label": label,
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished", tool_name=tool_name, ok=True, output=citation[:200]
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": True,
                "message": citation,
            }
        except Exception as e:
            yield {
                "type": "activity",
                "id": tool_id,
                "label": f"Failed: {label}",
                "icon": "file",
                "done": True,
            }
            if run is not None:
                run.log(
                    "tool_finished",
                    tool_name=tool_name,
                    ok=False,
                    output=_friendly_error(e),
                )
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": False,
                "message": _friendly_error(e),
            }
        return

    if run is not None:
        run.log(
            "tool_finished",
            tool_name=tool_name,
            ok=False,
            output=f"Unknown tool: {tool_name}",
        )
    yield {
        "type": "tool_result",
        "tool": tool_name,
        "ok": False,
        "message": f"Unknown tool: {tool_name}. This tool is not available — try a different approach.",
    }


# ---------------- Impls ----------------


def _short_path(path: str) -> str:
    try:
        home = str(Path.home())
        s = str(path)
        if s.startswith(home):
            s = "~" + s[len(home) :]
        return s
    except Exception:
        return str(path)


def _icon_for_path(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in (".html", ".htm"):
        return "html"
    if ext in (".css",):
        return "css"
    if ext in (".js", ".mjs"):
        return "js"
    if ext in (".ts", ".tsx"):
        return "ts"
    if ext in (".jsx",):
        return "jsx"
    if ext in (".json",):
        return "json"
    if ext in (".py",):
        return "code"
    if ext in (".md", ".markdown"):
        return "file"
    return "file"


def _write_file(path: str, content: str) -> None:
    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")


def _read_file(path: str) -> str:
    p = Path(path).expanduser()
    data = p.read_text(encoding="utf-8", errors="replace")
    # Cap to avoid flooding the context
    if len(data) > 200_000:
        return data[:200_000] + "\n…[truncated]"
    return data


def _list_dir(path: str) -> str:
    p = Path(path).expanduser()
    if not p.exists():
        raise FileNotFoundError(f"No such directory: {path}")
    if not p.is_dir():
        raise NotADirectoryError(f"Not a directory: {path}")
    entries = []
    for child in sorted(p.iterdir()):
        suffix = "/" if child.is_dir() else ""
        entries.append(child.name + suffix)
    return "\n".join(entries) if entries else "(empty)"


def _format_mcp_result(result: dict) -> str:
    parts: list[str] = []
    for item in result.get("content") or []:
        if not isinstance(item, dict):
            continue
        if item.get("type") == "text":
            parts.append(str(item.get("text") or ""))
        else:
            parts.append(json.dumps(item, ensure_ascii=False))
    text = "\n".join(p for p in parts if p).strip()
    return text or ("error" if result.get("isError") else "ok")


def _web_search(query: str, max_results: int = 6) -> str:
    max_results = max(1, min(int(max_results or 6), 10))
    base = "https://news.google.com/rss"
    params = {"hl": "en-US", "gl": "US", "ceid": "US:en"}
    if query:
        base += "/search"
        params["q"] = query
    url = base + "?" + urllib.parse.urlencode(params)
    with httpx.Client(timeout=12.0, follow_redirects=True) as client:
        resp = client.get(
            url,
            headers={
                "User-Agent": "zWork/1.0 (+https://tryzwork.app)",
                "Accept": "application/rss+xml, application/xml;q=0.9, */*;q=0.8",
            },
        )
        resp.raise_for_status()
        data = resp.content[:2_000_000]
    root = ET.fromstring(data)
    rows: list[str] = []
    for item in root.findall("./channel/item")[:max_results]:
        title = (item.findtext("title") or "").strip()
        link = (item.findtext("link") or "").strip()
        pub_date = (item.findtext("pubDate") or "").strip()
        source = ""
        for child in list(item):
            if child.tag.endswith("source"):
                source = (child.text or "").strip()
                break
        if not title:
            continue
        meta = " | ".join(p for p in (source, pub_date) if p)
        if meta:
            rows.append(f"- {title}\n  {meta}\n  {link}")
        else:
            rows.append(f"- {title}\n  {link}")
    if not rows:
        return "No web/news results found."
    heading = f"Results for: {query}" if query else "Top current headlines"
    return heading + "\n\n" + "\n".join(rows)


def _shell_path() -> str | None:
    """Return a working shell binary, falling back if /bin/sh is broken.

    On Windows returns None so subprocess uses the default COMSPEC.

    Tests each candidate with a command that exercises fork+exec
    (not just a shell builtin like ``true``), so broken shared library
    symbols are caught early and a working shell is selected.
    The result is cached — tested once per process lifetime.
    """
    if os.name == "nt":
        return None
    cached = getattr(_shell_path, "_cached", None)
    if cached is not None:
        return cached

    # Clean environment for shell testing: remove bundled libs that break system binaries
    test_env = os.environ.copy()
    for key in ("LD_LIBRARY_PATH", "GTK_PATH", "QT_PLUGIN_PATH", "GST_PLUGIN_PATH", "GST_PLUGIN_SYSTEM_PATH", "GST_PLUGIN_SYSTEM_PATH_1_0"):
        test_env.pop(key, None)
    for key, val in list(test_env.items()):
        if "/tmp/_MEI" in val or "/extracted/usr/lib" in val or "/extracted/lib" in val:
            test_env.pop(key, None)

    for candidate in (
        "/usr/bin/bash",
        "/bin/bash",
        "/bin/sh",
        "/usr/bin/sh",
    ):
        try:
            result = subprocess.run(
                [candidate, "-c", "echo ok && command -v cat >/dev/null"],
                capture_output=True,
                timeout=5,
                env=test_env,
            )
            if result.returncode == 0:
                _shell_path._cached = candidate
                return candidate
        except Exception:
            continue
    _shell_path._cached = "/bin/sh"
    return "/bin/sh"


async def _run_command(command: str, cwd: str) -> dict[str, Any]:
    _ensure_command_allowed(command)
    run = current_run()
    timeout_seconds = run.command_timeout_seconds if run is not None else 120
    output_cap = run.command_output_cap if run is not None else 20_000

    # Clean environment for subprocess: remove bundled PyInstaller libs that
    # break system binaries (e.g., bundled readline causes /bin/sh to fail)
    clean_env = os.environ.copy()
    for key in ("LD_LIBRARY_PATH", "GTK_PATH", "QT_PLUGIN_PATH", "GST_PLUGIN_PATH", "GST_PLUGIN_SYSTEM_PATH", "GST_PLUGIN_SYSTEM_PATH_1_0"):
        clean_env.pop(key, None)
    for key, val in list(clean_env.items()):
        if "/tmp/_MEI" in val or "/extracted/usr/lib" in val or "/extracted/lib" in val:
            clean_env.pop(key, None)

    kwargs: dict[str, Any] = {
        "cwd": str(Path(cwd).expanduser()),
        "stdout": asyncio.subprocess.PIPE,
        "stderr": asyncio.subprocess.PIPE,
        "env": clean_env,
    }
    if os.name != "nt":
        kwargs["start_new_session"] = True
        kwargs["executable"] = _shell_path()
    proc = await asyncio.create_subprocess_shell(command, **kwargs)
    if run is not None:
        run.register_process(proc.pid)
        run.log(
            "process_spawned",
            pid=proc.pid,
            command=command,
            cwd=str(Path(cwd).expanduser()),
        )
    try:
        try:
            stdout, stderr = await asyncio.wait_for(
                proc.communicate(), timeout=timeout_seconds
            )
        except asyncio.TimeoutError:
            _terminate_process_tree(proc.pid)
            await proc.wait()
            return {
                "ok": False,
                "returncode": -1,
                "output": f"Command timed out after {int(timeout_seconds)}s. Try a shorter command or break the work into smaller steps.",
            }
        output = (stdout or b"").decode("utf-8", errors="replace")
        stderr_text = (stderr or b"").decode("utf-8", errors="replace")
        if stderr_text:
            output += ("\n" + stderr_text) if output else stderr_text
        if len(output) > output_cap:
            output = output[:output_cap] + "\n…[truncated]"
        return {
            "ok": proc.returncode == 0,
            "returncode": proc.returncode,
            "output": output.strip(),
        }
    except asyncio.CancelledError:
        _terminate_process_tree(proc.pid)
        with contextlib.suppress(ProcessLookupError):
            await proc.wait()
        raise
    finally:
        if run is not None:
            run.unregister_process(proc.pid)


def _run_background(command: str, cwd: str) -> int:
    """Start a detached background process. Returns PID."""
    _ensure_command_allowed(command)

    # Clean environment for subprocess: remove bundled PyInstaller libs
    clean_env = os.environ.copy()
    for key in ("LD_LIBRARY_PATH", "GTK_PATH", "QT_PLUGIN_PATH", "GST_PLUGIN_PATH", "GST_PLUGIN_SYSTEM_PATH", "GST_PLUGIN_SYSTEM_PATH_1_0"):
        clean_env.pop(key, None)
    for key, val in list(clean_env.items()):
        if "/tmp/_MEI" in val or "/extracted/usr/lib" in val or "/extracted/lib" in val:
            clean_env.pop(key, None)

    kwargs: dict[str, Any] = {
        "shell": True,
        "cwd": Path(cwd).expanduser(),
        "stdout": subprocess.DEVNULL,
        "stderr": subprocess.DEVNULL,
        "stdin": subprocess.DEVNULL,
        "env": clean_env,
    }
    if os.name == "nt":
        # Windows: create a new process group so we can signal the tree
        kwargs["creationflags"] = subprocess.CREATE_NEW_PROCESS_GROUP  # type: ignore[assignment]
    else:
        kwargs["start_new_session"] = True
    proc = subprocess.Popen(command, **kwargs)
    return proc.pid


def _ensure_command_allowed(command: str) -> None:
    if _targets_zwork_backend(command):
        raise PermissionError(
            "Refusing to run a command that kills the zWork local backend on port 8787. "
            "Restart or inspect the backend instead of killing the app's own service."
        )


def _terminate_process_tree(pid: int) -> None:
    if pid <= 0:
        return
    if os.name == "nt":
        # Windows: use taskkill /T to terminate the process tree
        with contextlib.suppress(Exception):
            subprocess.run(
                ["taskkill", "/T", "/PID", str(pid), "/F"],
                capture_output=True,
                timeout=10,
            )
        return
    with contextlib.suppress(ProcessLookupError):
        os.killpg(pid, signal.SIGTERM)


def _pick_free_port(preferred: list[int]) -> int:
    import socket

    for port in preferred:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
            try:
                sock.bind(("127.0.0.1", port))
                return port
            except OSError:
                continue
    # fall back to kernel-assigned
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        sock.bind(("127.0.0.1", 0))
        return sock.getsockname()[1]


def _wait_for_port(port: int, timeout: float = 5.0) -> bool:
    """Poll localhost:port until it accepts connections or timeout expires."""
    import socket
    import time as _time

    deadline = _time.monotonic() + timeout
    while _time.monotonic() < deadline:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
            sock.settimeout(0.5)
            try:
                sock.connect(("127.0.0.1", port))
                return True
            except (ConnectionRefusedError, OSError):
                _time.sleep(0.25)
    return False


def _deploy_web_app(project_path: str, framework: str) -> dict[str, Any]:
    """Actually start a server."""
    p = Path(project_path).expanduser().resolve()
    if not p.exists():
        return {"ok": False, "message": f"Project path does not exist: {project_path}"}
    if not p.is_dir():
        return {
            "ok": False,
            "message": f"Project path is not a directory: {project_path}",
        }

    pkg = p / "package.json"
    if pkg.exists():
        try:
            data = json.loads(pkg.read_text())
            scripts = data.get("scripts") or {}
            if "dev" in scripts:
                port = _pick_free_port([5173, 3000, 8080])
                _run_background(f"PORT={port} npm run dev", str(p))
                if _wait_for_port(port):
                    return {
                        "ok": True,
                        "message": f"Started `npm run dev` in {p.name}. Open http://localhost:{port}",
                    }
                return {
                    "ok": False,
                    "message": f"`npm run dev` started but port {port} never opened. "
                    "The dev server may have crashed — check its output.",
                }
            if "start" in scripts:
                port = _pick_free_port([3000, 8080])
                _run_background(f"PORT={port} npm start", str(p))
                if _wait_for_port(port):
                    return {
                        "ok": True,
                        "message": f"Started `npm start` in {p.name}. Open http://localhost:{port}.",
                    }
                return {
                    "ok": False,
                    "message": f"`npm start` started but port {port} never opened. "
                    "The server may have crashed — check its output.",
                }
        except Exception:
            pass

    index = p / "index.html"
    if index.exists():
        port = _pick_free_port([8000, 8080, 5173])
        _run_background(f"python3 -m http.server {port}", str(p))
        if _wait_for_port(port):
            return {
                "ok": True,
                "message": f"Serving {p.name} at http://localhost:{port}",
            }
        return {
            "ok": False,
            "message": f"http.server started but port {port} never opened. "
            "The process may have crashed immediately.",
        }

    return {
        "ok": False,
        "message": f"No index.html or package.json in {p}. Nothing obvious to serve.",
    }


def _save_memory(content: str) -> None:
    """Append a note to the global memory file."""
    import time as _time

    p = memory_path()
    p.parent.mkdir(parents=True, exist_ok=True)
    existing = ""
    if p.exists():
        existing = p.read_text(encoding="utf-8").rstrip()
    timestamp = _time.strftime("%Y-%m-%d")
    entry = f"\n- {content}  ({timestamp})"
    p.write_text((existing + entry + "\n"), encoding="utf-8")


# ---------------- Document extraction ----------------

# Cap text payload to keep model context sane. Mirrors _read_file's 200k limit.
_EXTRACT_TEXT_CAP = 200_000

_EXTRACT_SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}


def _parse_page_range(spec: str, total: int) -> list[int]:
    """Parse a 1-based page spec like '1-5' or '3' into 0-based indices.

    Out-of-range pages are dropped silently rather than erroring, but a
    completely empty result (or a malformed spec) raises ValueError so the
    caller sees something actionable.
    """
    s = (spec or "").strip()
    if not s:
        raise ValueError("Empty page range")
    if "-" in s:
        parts = s.split("-", 1)
        try:
            start = int(parts[0])
            end = int(parts[1])
        except ValueError as e:
            raise ValueError(f"Invalid page range '{spec}'") from e
        if start < 1 or end < start:
            raise ValueError(f"Invalid page range '{spec}'")
    else:
        try:
            start = end = int(s)
        except ValueError as e:
            raise ValueError(f"Invalid page range '{spec}'") from e
        if start < 1:
            raise ValueError(f"Invalid page range '{spec}'")
    indices = [i - 1 for i in range(start, end + 1) if 1 <= i <= total]
    if not indices:
        raise ValueError(f"Page range '{spec}' is outside the document (1..{total})")
    return indices


def _cap_text(text: str, metadata: dict[str, Any]) -> str:
    if len(text) > _EXTRACT_TEXT_CAP:
        metadata["truncated"] = True
        return text[:_EXTRACT_TEXT_CAP]
    return text


_MAX_PDF_FILE_MB = 100
_MAX_PDF_PAGES = 500
_MAX_TABLE_PAGES = 50


def _extract_pdf(p: Path, fmt: str, pages: str | None) -> dict[str, Any]:
    import pypdf

    fsize_mb = p.stat().st_size / (1024 * 1024)
    if fsize_mb > _MAX_PDF_FILE_MB:
        raise ValueError(
            f"PDF is {fsize_mb:.0f} MB, exceeds the {_MAX_PDF_FILE_MB} MB limit. "
            "Split the file or extract specific pages."
        )

    reader = pypdf.PdfReader(str(p))
    total = len(reader.pages)
    if total > _MAX_PDF_PAGES:
        raise ValueError(
            f"PDF has {total} pages, exceeds the {_MAX_PDF_PAGES} page limit. "
            "Use the 'pages' parameter to extract a subset (e.g. '1-50')."
        )
    indices = _parse_page_range(pages, total) if pages else list(range(total))

    metadata: dict[str, Any] = {}
    info = reader.metadata or {}
    for key in ("/Title", "/Author", "/Subject", "/Creator", "/Producer"):
        val = info.get(key)
        if val:
            metadata[key.lstrip("/").lower()] = str(val)

    chunks: list[str] = []
    for idx in indices:
        try:
            page_text = reader.pages[idx].extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            chunks.append(page_text)

    text = "\n\n".join(chunks).strip()
    if not text:
        # No extractable text typically means a scanned image PDF. Surface
        # this rather than silently returning empty so the agent can decide
        # to fall back to OCR or tell the user.
        metadata["likely_scanned"] = True

    tables: list[dict[str, Any]] = []
    # pdfplumber is heavier than pypdf, only spin it up when tables are
    # actually requested (markdown output) and the doc isn't likely scanned.
    # Cap table extraction at _MAX_TABLE_PAGES to avoid OOM on huge PDFs.
    if fmt == "markdown" and not metadata.get("likely_scanned"):
        if len(indices) <= _MAX_TABLE_PAGES:
            try:
                import pdfplumber

                with pdfplumber.open(str(p)) as pdf:
                    for idx in indices:
                        if idx >= len(pdf.pages):
                            continue
                        for raw in pdf.pages[idx].extract_tables() or []:
                            if not raw:
                                continue
                            tables.append(
                                {
                                    "page": idx + 1,
                                    "rows": [
                                        [("" if c is None else str(c)) for c in row]
                                        for row in raw
                                    ],
                                }
                            )
            except Exception:
                # Table extraction is best-effort; never fail the whole call
                # because pdfplumber choked on a malformed table.
                pass
        else:
            metadata["tables_skipped"] = (
                f"Table extraction skipped: {len(indices)} pages exceeds "
                f"the {_MAX_TABLE_PAGES}-page limit for table extraction."
            )

    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": tables,
        "page_count": total,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_docx(p: Path, fmt: str) -> dict[str, Any]:
    import docx

    doc = docx.Document(str(p))
    paras = [para.text for para in doc.paragraphs if para.text]
    text = "\n\n".join(paras).strip()

    metadata: dict[str, Any] = {}
    core = doc.core_properties
    if core.title:
        metadata["title"] = core.title
    if core.author:
        metadata["author"] = core.author

    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": 0,  # docx has no fixed page count without rendering
        "metadata": metadata,
        "format": fmt,
    }


def _extract_xlsx(p: Path, fmt: str) -> dict[str, Any]:
    import openpyxl

    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    sheet_count = len(wb.sheetnames)
    text_chunks: list[str] = []
    tables: list[dict[str, Any]] = []
    for sheet in wb.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            rows.append([("" if v is None else str(v)) for v in row])
        # Strip trailing all-empty rows that openpyxl can leave behind.
        while rows and not any(c.strip() for c in rows[-1]):
            rows.pop()
        if not rows:
            continue
        tables.append({"sheet": sheet.title, "rows": rows})
        text_chunks.append(f"# {sheet.title}\n" + "\n".join("\t".join(r) for r in rows))
    wb.close()

    metadata: dict[str, Any] = {"sheet_count": sheet_count}
    text = "\n\n".join(text_chunks).strip()
    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": tables,
        "page_count": sheet_count,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_pptx(p: Path, fmt: str) -> dict[str, Any]:
    import pptx

    prs = pptx.Presentation(str(p))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    slide_lines.append(line)
        if slide_lines:
            chunks.append(f"# Slide {i}\n" + "\n".join(slide_lines))

    metadata: dict[str, Any] = {}
    text = "\n\n".join(chunks).strip()
    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": len(prs.slides),
        "metadata": metadata,
        "format": fmt,
    }


def _extract_plaintext(p: Path, fmt: str) -> dict[str, Any]:
    raw = p.read_text(encoding="utf-8", errors="replace")
    metadata: dict[str, Any] = {}
    text = _cap_text(raw, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": 0,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_document(path: str, fmt: str, pages: str | None) -> dict[str, Any]:
    p = Path(path).expanduser()
    if not p.exists():
        raise ValueError(f"File not found: {path}")
    if not p.is_file():
        raise ValueError(f"Not a file: {path}")

    ext = p.suffix.lower()
    if ext not in _EXTRACT_SUPPORTED:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(_EXTRACT_SUPPORTED))}"
        )

    if pages and ext != ".pdf":
        # Other formats don't have a stable page concept here. Better to
        # tell the caller than to silently ignore the filter.
        raise ValueError(f"'pages' is only supported for PDF, not {ext}")

    if ext == ".pdf":
        return _extract_pdf(p, fmt, pages)
    if ext == ".docx":
        return _extract_docx(p, fmt)
    if ext == ".xlsx":
        return _extract_xlsx(p, fmt)
    if ext == ".pptx":
        return _extract_pptx(p, fmt)
    return _extract_plaintext(p, fmt)


def _dctl_env() -> dict[str, str]:
    """Return a clean environment for running dctl as a standalone binary."""
    env = os.environ.copy()
    # Clean bundled PyInstaller libs that break system binaries
    for key in ("LD_LIBRARY_PATH", "PYTHONHOME", "PYTHONPATH", "GTK_PATH", "QT_PLUGIN_PATH", "GST_PLUGIN_PATH", "GST_PLUGIN_SYSTEM_PATH", "GST_PLUGIN_SYSTEM_PATH_1_0"):
        env.pop(key, None)
    for key, val in list(env.items()):
        if "/tmp/_MEI" in val or "/extracted/usr/lib" in val or "/extracted/lib" in val:
            env.pop(key, None)
    return env


def _dctl_path() -> str:
    """Find the dctl standalone binary in dev or bundle layouts."""
    this_file = Path(__file__).resolve()
    candidates = []

    # Bundled: same directory as the backend binary
    if "extracted" in str(this_file):
        bundle_dir = this_file.parents[2]  # up from .../extracted/usr/lib/python*/sidecar/agent/tools.py
        candidates.extend([
            bundle_dir / "usr" / "bin" / "dctl",
            bundle_dir / "bin" / "dctl",
            bundle_dir / "dctl",
        ])

    # Dev: zWork repo root
    for parent in this_file.parents:
        if (parent / "zWork-Skills").exists() or (parent / "sidecar").exists():
            candidates.append(parent.parent / "dctl" / "dctl")
            candidates.append(parent.parent / "dctl")
            break

    # Installed via install.sh
    candidates.append(Path.home() / ".local" / "bin" / "dctl")

    # System PATH
    import shutil
    system_dctl = shutil.which("dctl")
    if system_dctl:
        candidates.append(Path(system_dctl))

    for candidate in candidates:
        if candidate.exists() and os.access(candidate, os.X_OK):
            return str(candidate)

    # Fallback: hope it's in PATH
    return "dctl"


def _map_dctl_system(params: dict[str, Any]) -> tuple[str, list[str]]:
    action = params.get("action", "")
    target = params.get("target")
    args = []
    if target:
        args.append(target)
    return action, args


def _map_dctl_ui(params: dict[str, Any]) -> tuple[str, list[str]]:
    action = params.get("action", "")
    args = []

    if action in ("element", "read", "focus", "click", "wait"):
        selector = params.get("selector", "")
        args.append(selector)
        if action == "click":
            button = params.get("button", "left")
            if button != "left":
                args.extend(["--button", button])
            if params.get("double"):
                args.append("--double")
    elif action == "type":
        text = params.get("text", "")
        args.append(text)
        selector = params.get("selector")
        if selector:
            args.extend(["--into", selector])
    elif action == "key":
        combo = params.get("combo", "")
        args.append(combo)
    elif action == "scroll":
        direction = params.get("direction", "down")
        args.append(direction)
        amount = params.get("amount")
        if amount is not None:
            args.extend(["--amount", str(amount)])
    elif action == "describe":
        x = params.get("x")
        y = params.get("y")
        if x is not None and y is not None:
            args.extend([str(x), str(y)])
    elif action == "clipboard":
        clip_action = params.get("clipboard_action", "read")
        args.append(clip_action)
        if clip_action == "write":
            text = params.get("text", "")
            args.append(text)

    return action, args


def _map_dctl_browser(params: dict[str, Any]) -> tuple[str, list[str]]:
    action = params.get("action", "")
    target = params.get("target")
    selector = params.get("selector")
    url = params.get("url")
    text = params.get("text")
    expression = params.get("expression")
    session = params.get("session")

    # All browser commands start with browser subcommand
    args = [action]

    # Subcommand specific positioning
    if action in ("open", "activate", "close", "dom", "ax", "text", "actions", "selection", "caret", "snapshot", "click", "click-action", "act", "type", "press", "eval", "send", "wait-url", "wait-selector"):
        if target:
            args.append(target)

    if action in ("click", "wait-selector", "selector"):
        if selector:
            args.append(selector)
    elif action in ("type", "dom", "ax", "text", "caret"):
        if selector:
            args.extend(["--selector", selector])

    if action == "open" and url:
        args.append(url)
    elif action == "start" and url:
        args.extend(["--url", url])

    if action in ("type", "press", "eval", "send", "batch") and text:
        args.append(text)

    if expression and action == "eval":
        args.append(expression)

    if session:
        args.extend(["--session", session])

    return "browser", args


def _map_dctl_office(params: dict[str, Any]) -> tuple[str, list[str]]:
    otype = params.get("type", "")
    action = params.get("action", "")
    path = params.get("path", "")
    text = params.get("text")
    index = params.get("index")
    sheet = params.get("sheet")
    cell = params.get("cell")
    value = params.get("value")
    find = params.get("find")
    replace = params.get("replace")

    if otype == "libreoffice":
        args = [action]
        if path:
            args.append(path)
        if text:
            args.append(text)
        return "libreoffice", args

    subcommand = "docx" if otype == "word" else "xlsx"
    args = [action, path]

    if otype == "word":
        if action in ("append", "insert-before", "set-paragraph"):
            if index is not None and action in ("insert-before", "set-paragraph"):
                args.append(str(index))
            if text:
                args.append(text)
        elif action == "replace":
            if find:
                args.append(find)
            if replace:
                args.append(replace)
    elif otype == "excel":
        if sheet:
            args.append(sheet)
        if action == "read":
            range_val = cell or text or "A1"
            args.append(range_val)
        elif action == "write-cell":
            if cell:
                args.append(cell)
            if value:
                args.append(value)
        elif action == "write-range":
            if cell:
                args.append(cell)
            if value:
                args.append(value)
        elif action in ("locate-cell", "fill-cell"):
            if find:
                args.extend(["--row-label", find])
            if replace:
                args.extend(["--column-label", replace])
            if action == "fill-cell" and value:
                args.extend(["--value", value])
        elif action == "fill-table":
            if value:
                args.append(value)

    return subcommand, args


async def _run_dctl(subcommand: str, args: list[str], cwd: str) -> dict[str, Any]:
    if not subcommand:
        raise ValueError("dctl requires a subcommand")
    dctl_bin = _dctl_path()
    cmd = [dctl_bin, subcommand, *args]
    run = current_run()
    timeout_seconds = run.command_timeout_seconds if run is not None else 120
    output_cap = run.command_output_cap if run is not None else 20_000
    kwargs: dict[str, Any] = {
        "cwd": str(Path(cwd).expanduser()),
        "stdout": asyncio.subprocess.PIPE,
        "stderr": asyncio.subprocess.PIPE,
        "env": _dctl_env(),
    }
    if os.name != "nt":
        kwargs["start_new_session"] = True
    proc = await asyncio.create_subprocess_exec(*cmd, **kwargs)
    if run is not None:
        run.register_process(proc.pid)
        run.log(
            "process_spawned",
            pid=proc.pid,
            command=" ".join(cmd),
            cwd=str(Path(cwd).expanduser()),
        )
    try:
        try:
            stdout, stderr = await asyncio.wait_for(
                proc.communicate(), timeout=timeout_seconds
            )
        except asyncio.TimeoutError:
            _terminate_process_tree(proc.pid)
            await proc.wait()
            return {
                "ok": False,
                "returncode": -1,
                "output": f"dctl command timed out after {int(timeout_seconds)}s. Try a simpler operation.",
            }
    except asyncio.CancelledError:
        _terminate_process_tree(proc.pid)
        with contextlib.suppress(ProcessLookupError):
            await proc.wait()
        raise
    finally:
        if run is not None:
            run.unregister_process(proc.pid)
    output = (stdout or b"").decode("utf-8", errors="replace")
    stderr_text = (stderr or b"").decode("utf-8", errors="replace")
    if stderr_text:
        output += ("\n" + stderr_text) if output else stderr_text
    if len(output) > output_cap:
        output = output[:output_cap] + "\n…[truncated]"
    return {
        "ok": proc.returncode == 0,
        "returncode": proc.returncode,
        "output": output.strip(),
    }


# ---------------- Legacy <<TOOL>> marker parser ----------------


def parse_tool_calls(text: str) -> list[dict[str, Any]]:
    """Parse <<TOOL>>...<</TOOL>> blocks from model output (fallback)."""
    calls: list[dict[str, Any]] = []
    pattern = r"<<TOOL>>([\s\S]*?)<</TOOL>>"
    for match in re.finditer(pattern, text):
        try:
            data = json.loads(match.group(1).strip())
            if isinstance(data, dict) and "tool" in data and "params" in data:
                calls.append(data)
        except (json.JSONDecodeError, KeyError):
            continue
    return calls
